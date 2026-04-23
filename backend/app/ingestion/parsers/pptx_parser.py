import io
from app.ingestion.chunker import RawChunk, chunk_markdown


def parse_pptx(file_bytes: bytes, doc_title: str) -> list[RawChunk]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    sections: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []

        # Title shape (python-pptx exposes slide.shapes.title)
        title_shape = slide.shapes.title
        slide_title = title_shape.text.strip() if title_shape and title_shape.has_text_frame else ""

        heading = f"## Slide {i}" + (f": {slide_title}" if slide_title else "")
        parts.append(heading)

        # Body shapes (everything except the title placeholder)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape is title_shape:
                continue
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"*Notes: {notes_text}*")

        if len(parts) > 1:  # has more than just the heading
            sections.append("\n\n".join(parts))
        elif slide_title:
            sections.append(heading)

    content = "\n\n".join(sections)
    return chunk_markdown(content, source_title=doc_title, source_path=doc_title)
